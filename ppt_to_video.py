import time
import os
import requests
from datetime import datetime
from pptx import Presentation

# =============================================================================
# CONFIGURAÇÃO — preencha antes de rodar
# =============================================================================

API_KEY    = "COLE_SUA_API_KEY_AQUI"       # HeyGen → Settings → API
AVATAR_ID  = "COLE_SEU_AVATAR_ID_AQUI"     # HeyGen → Avatars → seu clone

# Caminho do PowerPoint (mesmo diretório do script)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ARQUIVO_PPT = os.path.join(SCRIPT_DIR, "apresentacao.pptx")

# Onde salvar o vídeo gerado (mesmo diretório, com timestamp)
timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
ARQUIVO_SAIDA = os.path.join(SCRIPT_DIR, f"video_{timestamp}.mp4")

# Estilo do vídeo — ajuste para replicar o vídeo de referência
VELOCIDADE_FALA = 1.0       # 0.5 (lento) a 2.0 (rápido)
LARGURA  = 1280             # resolução horizontal
ALTURA   = 720              # resolução vertical (use 1920x1080 para full HD)

# Cor de fundo (hex) — ou troque por "background_url" com imagem
COR_FUNDO = "#FFFFFF"

# =============================================================================

def extrair_texto_do_pptx(caminho):
    """Lê cada slide e retorna lista de textos."""
    if not os.path.exists(caminho):
        raise FileNotFoundError(f"Arquivo não encontrado: {caminho}")

    prs = Presentation(caminho)
    textos = []
    for i, slide in enumerate(prs.slides):
        partes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texto = shape.text.strip()
                if texto:
                    partes.append(texto)
        if partes:
            texto_slide = " ".join(partes)
            textos.append(texto_slide)
            print(f"  Slide {i+1}: {texto_slide[:70]}...")
    return textos


def gerar_video(script):
    """Envia o script para HeyGen e retorna o video_id."""
    print("\nEnviando para HeyGen API...")
    response = requests.post(
        "https://api.heygen.com/v2/video/generate",
        headers={
            "X-Api-Key": API_KEY,
            "Content-Type": "application/json"
        },
        json={
            "video_inputs": [{
                "character": {
                    "type": "avatar",
                    "avatar_id": AVATAR_ID,
                    "avatar_style": "normal"  # "normal" | "circle" | "closeUp"
                },
                "voice": {
                    "type": "text",
                    "input_text": script,
                    "speed": VELOCIDADE_FALA
                },
                "background": {
                    "type": "color",
                    "value": COR_FUNDO
                }
            }],
            "dimension": {
                "width": LARGURA,
                "height": ALTURA
            }
        }
    )

    if response.status_code != 200:
        print(f"Erro na requisição: {response.status_code} — {response.text}")
        return None

    data = response.json()
    video_id = data.get("data", {}).get("video_id")
    if video_id:
        print(f"Vídeo enfileirado com sucesso. video_id: {video_id}")
    else:
        print("Erro: video_id não retornado.", data)
    return video_id


def aguardar_e_baixar(video_id, caminho_saida):
    """Aguarda o processamento e baixa o vídeo ao concluir."""
    print(f"\nAguardando geração (pode levar 2 a 10 minutos)...")
    tentativa = 0
    while True:
        tentativa += 1
        response = requests.get(
            f"https://api.heygen.com/v1/video_status.get?video_id={video_id}",
            headers={"X-Api-Key": API_KEY}
        )
        status = response.json().get("data", {})
        estado = status.get("status", "desconhecido")
        print(f"  [{tentativa}] Status: {estado}")

        if estado == "completed":
            url = status.get("video_url")
            print(f"\nBaixando vídeo de: {url}")
            video = requests.get(url, stream=True)
            with open(caminho_saida, "wb") as f:
                for chunk in video.iter_content(chunk_size=8192):
                    f.write(chunk)
            print(f"\nVídeo salvo em:\n  {caminho_saida}")
            break

        elif estado == "failed":
            erro = status.get("error", {})
            print(f"\nFalha na geração: {erro}")
            break

        time.sleep(20)


# =============================================================================
# Execução principal
# =============================================================================

print("=== PPT to Video — HeyGen API ===\n")
print(f"Lendo apresentação: {ARQUIVO_PPT}")
slides = extrair_texto_do_pptx(ARQUIVO_PPT)
print(f"\n{len(slides)} slide(s) encontrado(s).")

# Junta todos os slides em um único script contínuo
script_completo = ". ".join(slides)
print(f"\nScript total: {len(script_completo)} caracteres")

video_id = gerar_video(script_completo)
if video_id:
    aguardar_e_baixar(video_id, ARQUIVO_SAIDA)
